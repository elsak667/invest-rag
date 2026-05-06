import sys
import math
import requests
import fitz
from pptx import Presentation
import pandas as pd
from pathlib import Path

# 配置
VPS_URL = "https://nottingham-protected-trivia-cassette.trycloudflare.com/embed"
SUPABASE_URL = "https://rgnncmgrumwjjgzyhmkt.supabase.co"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6InJnbm5jbWdydW13ampnenlobWt0Iiwicm9sZSI6InNlcnZpY2Vfcm9sZSIsImlhdCI6MTc3NzczMTAwNSwiZXhwIjoyMDkzMzA3MDA1fQ.Rrq95ziEwcTmX1NNz62ZqGsN8GvyZcutqzAIwC7PkC8"
CHUNK_SIZE = 500

HEADERS = {
    "apikey": SUPABASE_KEY,
    "Authorization": f"Bearer {SUPABASE_KEY}",
    "Content-Type": "application/json",
    "Prefer": "return=minimal"
}

def embed_texts(texts):
    resp = requests.post(VPS_URL, json={"texts": texts}, timeout=60)
    resp.raise_for_status()
    return resp.json()["embeddings"]

def extract_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_pdf(path):
    doc = fitz.open(path)
    texts = []
    for page in doc:
        t = page.get_text()
        if t.strip():
            texts.append(t)
    return "\n".join(texts)

def extract_excel(path):
    dfs = pd.read_excel(path, sheet_name=None)
    texts = []
    for sheet, df in dfs.items():
        texts.append(f"Sheet: {sheet}\n{df.to_string()}")
    return "\n".join(texts)

def chunk_text(text, size=CHUNK_SIZE):
    words = text.split()
    chunk, current = [], []
    for w in words:
        current.append(w)
        if sum(len(x) for x in current) >= size:
            chunk.append(" ".join(current))
            current = []
    if current:
        chunk.append(" ".join(current))
    return chunk

def sanitize_vector(emb):
    result = []
    for v in emb:
        try:
            f = float(v)
            if math.isnan(f) or math.isinf(f):
                f = 0.0
            result.append(f)
        except (ValueError, TypeError):
            result.append(0.0)
    return result

def get_session():
    s = requests.Session()
    s.trust_env = False  # ignore system proxy
    return s

_session = get_session()

def store_doc(company, doc_type, source, content, embedding):
    # 清洗 content 中的空字符
    content = content.replace('\x00', '')
    payload = {
        "company": company,
        "doc_type": doc_type,
        "source": source,
        "content": content,
        "embedding": sanitize_vector(embedding)
    }
    resp = _session.post(
        f"{SUPABASE_URL}/rest/v1/documents",
        headers=HEADERS,
        json=payload
    )
    if resp.status_code not in (200, 201):
        print(f"  ERROR {resp.status_code}: {resp.text[:200]}")
        resp.raise_for_status()

def process_file(path, company, doc_type):
    ext = Path(path).suffix.lower()
    print(f"Processing {path} ({ext})")
    if ext == ".pptx":
        text = extract_pptx(path)
    elif ext == ".pdf":
        text = extract_pdf(path)
    elif ext in [".xlsx", ".xls"]:
        text = extract_excel(path)
    else:
        print(f"Unsupported: {ext}")
        return

    chunks = chunk_text(text)
    print(f"  {len(chunks)} chunks extracted")

    # batch embedding
    embeddings = embed_texts(chunks)
    print(f"  embedding done, storing...")

    for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
        try:
            store_doc(company, doc_type, path, chunk, emb)
        except Exception as e:
            print(f"  chunk {i+1} failed: {e}")
            raise
        if (i + 1) % 10 == 0:
            print(f"  stored {i+1}/{len(chunks)}")

    print(f"  ✓ {len(chunks)} chunks stored in Supabase")

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python3 embed_doc.py <file> <company> <doc_type>")
        sys.exit(1)
    process_file(sys.argv[1], sys.argv[2], sys.argv[3])
